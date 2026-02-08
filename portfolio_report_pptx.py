"""
portfolio_report_pptx.py
────────────────────────
Professional PowerPoint portfolio report generator.

Usage:
    from portfolio_report_pptx import generate_pptx_report
    generate_pptx_report(report_data, "portfolio_report.pptx")

Dependencies: python-pptx, matplotlib, pandas, numpy, Pillow
"""

from __future__ import annotations

import io
import os
import tempfile
from datetime import datetime
from typing import Any, Dict, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# Color palette
# ──────────────────────────────────────────────
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
C_NAVY      = RGBColor(0x0F, 0x34, 0x60)
C_ACCENT    = RGBColor(0xE9, 0x45, 0x60)
C_BLUE      = RGBColor(0x00, 0xB4, 0xD8)
C_GOLD      = RGBColor(0xFC, 0xA3, 0x11)
C_GREEN     = RGBColor(0x6A, 0x99, 0x4E)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF0, 0xF0, 0xF5)
C_MUTED     = RGBColor(0x66, 0x66, 0x66)
C_BODY      = RGBColor(0x33, 0x33, 0x33)

PIE_COLORS_RGB = [
    (0x0F, 0x34, 0x60), (0xE9, 0x45, 0x60), (0x00, 0xB4, 0xD8), (0xFC, 0xA3, 0x11),
    (0x6A, 0x99, 0x4E), (0x9B, 0x5D, 0xE5), (0xF1, 0x5B, 0xB5), (0x00, 0xF5, 0xD4),
    (0xFB, 0x85, 0x00), (0x21, 0x9E, 0xBC), (0x83, 0x38, 0xEC), (0xFF, 0x00, 0x6E),
]

PIE_COLORS_HEX = [
    "#0f3460", "#e94560", "#00b4d8", "#fca311",
    "#6a994e", "#9b5de5", "#f15bb5", "#00f5d4",
    "#fb8500", "#219ebc", "#8338ec", "#ff006e",
]


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────
def _set_cell_text(cell, text: str, font_size: int = 10,
                   bold: bool = False, color: RGBColor = C_BODY,
                   alignment: PP_ALIGN = PP_ALIGN.LEFT):
    """Set text in a table cell with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_shape_bg(slide, left, top, width, height, fill_color: RGBColor):
    """Add a filled rectangle as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size: int = 12, bold: bool = False,
                 color: RGBColor = C_BODY, alignment: PP_ALIGN = PP_ALIGN.LEFT,
                 font_name: str = "Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def _fig_to_image_stream(fig: plt.Figure, dpi: int = 200) -> io.BytesIO:
    """Save matplotlib figure to a BytesIO PNG stream."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                pad_inches=0.1, facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ──────────────────────────────────────────────
# Chart builders (matplotlib → PNG)
# ──────────────────────────────────────────────
def _make_donut_chart(weights: Dict[str, float]) -> io.BytesIO:
    """Create a professional donut chart."""
    labels = list(weights.keys())
    sizes = [abs(v) for v in weights.values()]

    # Filter negligible
    filtered = [(l, s) for l, s in zip(labels, sizes) if s > 0.001]
    if not filtered:
        fig, ax = plt.subplots(figsize=(7, 5))
        ax.text(0.5, 0.5, "No allocations", ha="center", va="center", fontsize=14)
        ax.axis("off")
        return _fig_to_image_stream(fig)

    labels, sizes = zip(*filtered)
    colors = PIE_COLORS_HEX[:len(labels)]

    fig, ax = plt.subplots(figsize=(8, 5))
    wedges, texts, autotexts = ax.pie(
        sizes,
        labels=None,
        autopct=lambda p: f"{p:.1f}%" if p > 3 else "",
        startangle=140,
        colors=colors,
        pctdistance=0.78,
        wedgeprops=dict(width=0.40, edgecolor="white", linewidth=2),
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color("white")
        at.set_fontweight("bold")

    ax.legend(
        wedges,
        [f"{l}  ({s:.1%})" for l, s in zip(labels, sizes)],
        loc="center left",
        bbox_to_anchor=(1.05, 0.5),
        fontsize=10,
        frameon=False,
        labelspacing=1.0,
    )
    ax.set_aspect("equal")
    fig.subplots_adjust(left=0.05, right=0.60)
    return _fig_to_image_stream(fig)


def _make_performance_chart(
    portfolio_series: Optional[pd.Series],
    benchmark_series: Optional[pd.Series],
    portfolio_label: str = "Portfolio",
    benchmark_label: str = "Benchmark",
) -> Optional[io.BytesIO]:
    """Create a performance line chart."""
    if portfolio_series is None or portfolio_series.empty:
        return None

    fig, ax = plt.subplots(figsize=(10, 4.5))

    ax.plot(portfolio_series.index, portfolio_series.values,
            color="#0f3460", linewidth=2.2, label=portfolio_label)

    if benchmark_series is not None and not benchmark_series.empty:
        ax.plot(benchmark_series.index, benchmark_series.values,
                color="#e94560", linewidth=1.8, linestyle="--", label=benchmark_label, alpha=0.8)

    ax.set_ylabel("Growth of 100", fontsize=11, color="#333333")
    ax.tick_params(labelsize=9, colors="#666666")
    ax.legend(fontsize=10, frameon=True, fancybox=True, shadow=False,
              edgecolor="#cccccc", loc="upper left")
    ax.grid(True, alpha=0.25, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    fig.tight_layout()
    return _fig_to_image_stream(fig)


def _make_drawdown_chart(portfolio_series: pd.Series) -> Optional[io.BytesIO]:
    """Create a drawdown chart."""
    if portfolio_series is None or portfolio_series.empty:
        return None

    running_max = portfolio_series.cummax()
    drawdown = (portfolio_series - running_max) / running_max

    fig, ax = plt.subplots(figsize=(10, 3))
    ax.fill_between(drawdown.index, drawdown.values, 0,
                    color="#e94560", alpha=0.35)
    ax.plot(drawdown.index, drawdown.values, color="#e94560", linewidth=1.2)
    ax.set_ylabel("Drawdown", fontsize=10, color="#333333")
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda y, _: f"{y:.0%}"))
    ax.tick_params(labelsize=9, colors="#666666")
    ax.grid(True, alpha=0.2, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    fig.tight_layout()
    return _fig_to_image_stream(fig)


def _make_frontier_chart(
    pv_cloud: np.ndarray,
    pr_cloud: np.ndarray,
    port_v: float,
    port_r: float,
    cml_v: Optional[list] = None,
    cml_r: Optional[list] = None,
    selected_label: str = "Selected",
) -> io.BytesIO:
    """Create an efficient frontier scatter chart."""
    fig, ax = plt.subplots(figsize=(10, 5.5))

    # Cloud
    ax.scatter(pv_cloud, pr_cloud, s=3, alpha=0.25, c="#a0b4d0", edgecolors="none",
               label="Simulated portfolios", rasterized=True)

    # CML line
    if cml_v and cml_r:
        ax.plot(cml_v, cml_r, color="#fca311", linewidth=2.2, label="CML", zorder=5)

    # Max Sharpe from cloud
    sharpe_cloud = np.where(np.array(pv_cloud) > 1e-12,
                            (np.array(pr_cloud)) / np.array(pv_cloud), np.nan)
    i_best = int(np.nanargmax(sharpe_cloud))
    ax.scatter([pv_cloud[i_best]], [pr_cloud[i_best]], s=160, marker="*",
               c="#fca311", edgecolors="#333", linewidths=0.6,
               label="Max Sharpe (cloud)", zorder=6)

    # Selected portfolio
    ax.scatter([port_v], [port_r], s=120, marker="D", c="#e94560",
               edgecolors="white", linewidths=1.2, label=selected_label, zorder=7)

    ax.set_xlabel("Annualized Volatility", fontsize=11, color="#333")
    ax.set_ylabel("Annualized Return", fontsize=11, color="#333")
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda y, _: f"{y:.0%}"))
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda y, _: f"{y:.0%}"))
    ax.tick_params(labelsize=9, colors="#666")
    ax.legend(fontsize=9, frameon=True, fancybox=True, shadow=False,
              edgecolor="#ccc", loc="upper left")
    ax.grid(True, alpha=0.2, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#ccc")
    ax.spines["bottom"].set_color("#ccc")
    fig.tight_layout()
    return _fig_to_image_stream(fig)


def _make_cumulative_returns_chart(
    cum_returns: pd.DataFrame,
) -> Optional[io.BytesIO]:
    """Create a cumulative returns line chart for all individual assets."""
    if cum_returns is None or cum_returns.empty:
        return None

    fig, ax = plt.subplots(figsize=(10, 5))
    for i, col in enumerate(cum_returns.columns):
        ax.plot(cum_returns.index, cum_returns[col],
                linewidth=1.5, label=col,
                color=PIE_COLORS_HEX[i % len(PIE_COLORS_HEX)])
    ax.set_ylabel("Growth of 100", fontsize=11, color="#333333")
    ax.set_title("Cumulative Returns (Base = 100)", fontsize=13,
                 fontweight="bold", color="#1a1a2e", pad=10)
    ax.legend(fontsize=8, frameon=True, edgecolor="#ccc", loc="upper left",
              ncol=max(1, len(cum_returns.columns) // 6))
    ax.tick_params(labelsize=9, colors="#666666")
    ax.grid(True, alpha=0.25, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#cccccc")
    ax.spines["bottom"].set_color("#cccccc")
    fig.autofmt_xdate(rotation=30, ha="right")
    fig.tight_layout()
    return _fig_to_image_stream(fig)


def _make_correlation_heatmap(
    corr_matrix: pd.DataFrame,
) -> Optional[io.BytesIO]:
    """Create a correlation matrix heatmap."""
    if corr_matrix is None or corr_matrix.empty:
        return None

    fig, ax = plt.subplots(figsize=(8, 6))
    n = len(corr_matrix)
    im = ax.imshow(corr_matrix.values, cmap="RdBu_r", vmin=-1, vmax=1, aspect="auto")
    ax.set_xticks(range(n))
    ax.set_yticks(range(n))
    ax.set_xticklabels(corr_matrix.columns, fontsize=8, rotation=45, ha="right")
    ax.set_yticklabels(corr_matrix.index, fontsize=8)
    # Annotate cells
    for i in range(n):
        for j in range(n):
            val = corr_matrix.iloc[i, j]
            color = "white" if abs(val) > 0.6 else "black"
            ax.text(j, i, f"{val:.2f}", ha="center", va="center",
                    fontsize=8, color=color)
    fig.colorbar(im, ax=ax, shrink=0.8)
    ax.set_title("Correlation Matrix", fontsize=13, fontweight="bold",
                 color="#1a1a2e", pad=10)
    fig.tight_layout()
    return _fig_to_image_stream(fig)


# ──────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────
def _slide_cover(prs: Presentation, report_data: dict):
    """Slide 1: Title / cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height

    # Dark background bar (top 45%)
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(3.4), C_DARK)

    # Accent stripe
    _add_shape_bg(slide, Emu(0), Inches(3.4), sw, Inches(0.06), C_ACCENT)

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.8),
                 "PORTFOLIO REPORT", font_size=32, bold=True, color=C_WHITE,
                 font_name="Calibri")

    # Subtitle
    client = report_data.get("client_name", "Client")
    _add_textbox(slide, Inches(0.8), Inches(1.7), Inches(8), Inches(0.5),
                 f"Prepared for {client}", font_size=16, color=RGBColor(0xBB, 0xBB, 0xCC),
                 font_name="Calibri")

    # Date + value
    report_date = report_data.get("report_date", datetime.today().strftime("%Y-%m-%d"))
    currency = report_data.get("currency", "USD")
    value = report_data.get("portfolio_value")
    info_text = f"Report Date: {report_date}"
    if value is not None:
        info_text += f"   |   Portfolio Value: {value:,.2f} {currency}"
    _add_textbox(slide, Inches(0.8), Inches(2.4), Inches(8), Inches(0.4),
                 info_text, font_size=12, color=RGBColor(0x99, 0x99, 0xAA),
                 font_name="Calibri")

    # Strategy
    strategy = report_data.get("portfolio_label", "")
    if strategy:
        _add_textbox(slide, Inches(0.8), Inches(4.0), Inches(8), Inches(0.5),
                     f"Strategy: {strategy}", font_size=14, bold=True, color=C_NAVY)

    # Measurement currency disclaimer
    currency_note = report_data.get("measurement_currency_note", "")
    if currency_note:
        _add_textbox(slide, Inches(0.8), Inches(4.6), Inches(8), Inches(0.4),
                     f"\ud83d\udcd0 {currency_note}", font_size=10, color=C_MUTED,
                     font_name="Calibri")

    # Footer branding
    _add_textbox(slide, Inches(0.8), Inches(6.8), Inches(8), Inches(0.3),
                 "Generated by Portfolio Optimizer", font_size=9, color=C_MUTED,
                 alignment=PP_ALIGN.LEFT)


def _slide_kpis(prs: Presentation, report_data: dict):
    """Slide 2: Key performance indicators in card layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header bar
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "KEY METRICS", font_size=22, bold=True, color=C_WHITE)

    # Build KPI data
    perf = report_data.get("performance", {})
    risk = report_data.get("risk", {})
    currency = report_data.get("currency", "USD")
    value = report_data.get("portfolio_value", 0)

    kpis = []

    # Performance KPIs
    for period, ret in perf.items():
        if isinstance(ret, (int, float)):
            kpis.append((str(period), f"{ret:.2%}",
                         C_GREEN if ret >= 0 else C_ACCENT))

    # Risk KPIs
    for metric, val in risk.items():
        if isinstance(val, (int, float)):
            if "sharpe" in metric.lower():
                kpis.append((metric, f"{val:.2f}", C_NAVY))
            elif "drawdown" in metric.lower():
                kpis.append((metric, f"{val:.2%}", C_ACCENT))
            else:
                kpis.append((metric, f"{val:.2%}", C_NAVY))

    # Expected annual return in cash
    exp_r = risk.get("Expected Annual Return")
    if exp_r is None and value and perf:
        # Use 1Y return or Since Inception annualized
        ann = perf.get("1Y", perf.get("Since Inception", 0))
        if isinstance(ann, (int, float)):
            kpis.append((f"Exp. Return ({currency})", f"{value * ann:+,.0f} {currency}",
                         C_GREEN if ann >= 0 else C_ACCENT))

    # Expected Future Value cards
    horizon = report_data.get("investment_horizon")
    fv = report_data.get("expected_future_value")
    if horizon is not None and fv is not None:
        currency = report_data.get("currency", "USD")
        pnl_h = report_data.get("expected_pnl_horizon", 0)
        kpis.append((f"Expected Value ({horizon}Y)", f"{fv:,.0f} {currency}",
                     C_GREEN if pnl_h >= 0 else C_ACCENT))
        kpis.append((f"Expected P&L ({horizon}Y)", f"{pnl_h:+,.0f} {currency}",
                     C_GREEN if pnl_h >= 0 else C_ACCENT))

    # Lay out KPI cards in a 2-row grid
    cols = 4
    card_w = Inches(2.1)
    card_h = Inches(1.6)
    gap_x = Inches(0.25)
    start_x = Inches(0.5)
    start_y = Inches(1.3)

    for i, (label, val_str, val_color) in enumerate(kpis[:8]):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + Inches(0.2))

        # Card background
        card = _add_shape_bg(slide, x, y, card_w, card_h, C_LIGHT_BG)
        card.shadow.inherit = False

        # Value
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.25), card_w - Inches(0.3), Inches(0.7),
                     val_str, font_size=22, bold=True, color=val_color,
                     alignment=PP_ALIGN.CENTER)
        # Label
        _add_textbox(slide, x + Inches(0.15), y + Inches(1.0), card_w - Inches(0.3), Inches(0.4),
                     label, font_size=10, color=C_MUTED, alignment=PP_ALIGN.CENTER)


def _slide_allocation(prs: Presentation, report_data: dict):
    """Slide 3: Asset allocation donut chart."""
    weights = report_data.get("weights")
    if not weights:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "ASSET ALLOCATION", font_size=22, bold=True, color=C_WHITE)

    # Donut chart
    img_stream = _make_donut_chart(weights)
    slide.shapes.add_picture(img_stream, Inches(0.6), Inches(1.2),
                             width=Inches(8.8), height=Inches(5.0))


def _slide_holdings(prs: Presentation, report_data: dict):
    """Slide 4: Holdings table."""
    holdings_df = report_data.get("holdings_df")
    if holdings_df is None or not isinstance(holdings_df, pd.DataFrame) or holdings_df.empty:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "HOLDINGS", font_size=22, bold=True, color=C_WHITE)

    # Table
    n_rows = min(len(holdings_df) + 1, 16)  # cap rows
    n_cols = len(holdings_df.columns)
    tbl_width = Inches(9.0)
    tbl_height = Inches(0.4) * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        tbl_width, tbl_height,
    )
    table = table_shape.table

    # Style header row
    for j, col_name in enumerate(holdings_df.columns):
        cell = table.cell(0, j)
        _set_cell_text(cell, str(col_name), font_size=10, bold=True, color=C_WHITE,
                       alignment=PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY

    # Data rows
    currency = report_data.get("currency", "USD")
    for i, (_, row) in enumerate(holdings_df.head(n_rows - 1).iterrows()):
        for j, col_name in enumerate(holdings_df.columns):
            cell = table.cell(i + 1, j)
            val = row[col_name]
            if isinstance(val, float):
                if "weight" in col_name.lower():
                    text = f"{val:.1%}"
                elif "p/l" in col_name.lower() or "p&l" in col_name.lower():
                    text = f"{val:+,.2f} {currency}"
                else:
                    text = f"{val:,.2f}"
            else:
                text = str(val)

            color = C_BODY
            if isinstance(val, (int, float)) and ("p/l" in col_name.lower() or "p&l" in col_name.lower()):
                color = C_GREEN if val >= 0 else C_ACCENT

            _set_cell_text(cell, text, font_size=9, color=color,
                           alignment=PP_ALIGN.CENTER)
            # Alternate row shading
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_BG
            else:
                cell.fill.background()


def _slide_performance(prs: Presentation, report_data: dict):
    """Slide 5: Performance chart + drawdown."""
    portfolio_series = report_data.get("portfolio_series")
    if portfolio_series is None:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "PERFORMANCE", font_size=22, bold=True, color=C_WHITE)

    # Performance line chart
    perf_img = _make_performance_chart(
        portfolio_series,
        report_data.get("benchmark_series"),
        report_data.get("portfolio_label", "Portfolio"),
        report_data.get("benchmark_label", "Benchmark"),
    )
    if perf_img:
        slide.shapes.add_picture(perf_img, Inches(0.3), Inches(1.1),
                                 width=Inches(9.3), height=Inches(3.8))

    # Drawdown chart (below)
    dd_img = _make_drawdown_chart(portfolio_series)
    if dd_img:
        _add_textbox(slide, Inches(0.5), Inches(4.85), Inches(3), Inches(0.3),
                     "DRAWDOWN", font_size=10, bold=True, color=C_MUTED)
        slide.shapes.add_picture(dd_img, Inches(0.3), Inches(5.1),
                                 width=Inches(9.3), height=Inches(2.1))


def _slide_frontier(prs: Presentation, report_data: dict):
    """Slide: Efficient frontier chart."""
    frontier = report_data.get("frontier")
    if not frontier:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "EFFICIENT FRONTIER", font_size=22, bold=True, color=C_WHITE)

    img_stream = _make_frontier_chart(
        pv_cloud=frontier["pv_cloud"],
        pr_cloud=frontier["pr_cloud"],
        port_v=frontier["port_v"],
        port_r=frontier["port_r"],
        cml_v=frontier.get("cml_v"),
        cml_r=frontier.get("cml_r"),
        selected_label=frontier.get("selected_label", "Selected"),
    )
    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.1),
                             width=Inches(9.3), height=Inches(5.5))


def _slide_cumulative_returns(prs: Presentation, report_data: dict):
    """Slide: Cumulative returns for all individual assets."""
    cum_returns = report_data.get("cum_returns")
    if cum_returns is None or (hasattr(cum_returns, "empty") and cum_returns.empty):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "CUMULATIVE RETURNS", font_size=22, bold=True, color=C_WHITE)

    img_stream = _make_cumulative_returns_chart(cum_returns)
    if img_stream:
        slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.1),
                                 width=Inches(9.3), height=Inches(5.5))


def _slide_correlation(prs: Presentation, report_data: dict):
    """Slide: Correlation matrix heatmap."""
    corr_matrix = report_data.get("corr_matrix")
    if corr_matrix is None or (hasattr(corr_matrix, "empty") and corr_matrix.empty):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "CORRELATION MATRIX", font_size=22, bold=True, color=C_WHITE)

    img_stream = _make_correlation_heatmap(corr_matrix)
    if img_stream:
        slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.1),
                                 width=Inches(9.3), height=Inches(5.5))


def _slide_betas_capm(prs: Presentation, report_data: dict):
    """Slide: Betas table and CAPM regression results."""
    betas = report_data.get("betas")
    capm_df = report_data.get("capm_df")
    bench_label = report_data.get("benchmark_label", "Benchmark")

    if betas is None and (capm_df is None or (hasattr(capm_df, "empty") and capm_df.empty)):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 f"BETAS & CAPM vs {bench_label}", font_size=22, bold=True, color=C_WHITE)

    y_pos = Inches(1.2)

    # Betas table
    if betas is not None and len(betas) > 0:
        sorted_betas = betas.sort_values(ascending=False)
        n_rows = min(len(sorted_betas) + 1, 15)
        n_cols = 2

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), y_pos,
            Inches(4.0), Inches(0.35) * n_rows,
        )
        table = table_shape.table

        _set_cell_text(table.cell(0, 0), "Asset", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
        _set_cell_text(table.cell(0, 1), f"Beta vs {bench_label}", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = C_NAVY
        table.cell(0, 1).fill.solid()
        table.cell(0, 1).fill.fore_color.rgb = C_NAVY

        for i, (asset, beta_val) in enumerate(sorted_betas.head(n_rows - 1).items()):
            _set_cell_text(table.cell(i + 1, 0), str(asset), font_size=9, alignment=PP_ALIGN.CENTER)
            _set_cell_text(table.cell(i + 1, 1), f"{beta_val:.2f}", font_size=9, alignment=PP_ALIGN.CENTER)
            if i % 2 == 0:
                table.cell(i + 1, 0).fill.solid()
                table.cell(i + 1, 0).fill.fore_color.rgb = C_LIGHT_BG
                table.cell(i + 1, 1).fill.solid()
                table.cell(i + 1, 1).fill.fore_color.rgb = C_LIGHT_BG

    # CAPM table on right side
    if capm_df is not None and isinstance(capm_df, pd.DataFrame) and not capm_df.empty:
        capm_display = capm_df.copy()
        if "Alpha (daily)" in capm_display.columns:
            capm_display["Alpha (ann.)"] = (1.0 + capm_display["Alpha (daily)"]) ** 252 - 1.0
            display_cols = ["Asset", "Alpha (ann.)", "Beta", "R^2"]
            display_cols = [c for c in display_cols if c in capm_display.columns]
        else:
            display_cols = list(capm_display.columns)[:5]

        n_rows_capm = min(len(capm_display) + 1, 15)
        n_cols_capm = len(display_cols)

        table_shape2 = slide.shapes.add_table(
            n_rows_capm, n_cols_capm,
            Inches(5.0), y_pos,
            Inches(4.8), Inches(0.35) * n_rows_capm,
        )
        table2 = table_shape2.table

        for j, col_name in enumerate(display_cols):
            _set_cell_text(table2.cell(0, j), col_name, font_size=9, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
            table2.cell(0, j).fill.solid()
            table2.cell(0, j).fill.fore_color.rgb = C_NAVY

        for i, (_, crow) in enumerate(capm_display.head(n_rows_capm - 1).iterrows()):
            for j, col_name in enumerate(display_cols):
                val = crow[col_name]
                if col_name == "Asset":
                    text = str(val)
                elif "alpha" in col_name.lower():
                    text = f"{val:.2%}"
                else:
                    text = f"{val:.2f}" if isinstance(val, (int, float)) else str(val)
                _set_cell_text(table2.cell(i + 1, j), text, font_size=9, alignment=PP_ALIGN.CENTER)
                if i % 2 == 0:
                    table2.cell(i + 1, j).fill.solid()
                    table2.cell(i + 1, j).fill.fore_color.rgb = C_LIGHT_BG


def _slide_risk_analysis(prs: Presentation, report_data: dict):
    """Slide: Risk analysis table + daily returns statistics."""
    risk = report_data.get("risk", {})
    portfolio_series = report_data.get("portfolio_series")

    if not risk and (portfolio_series is None or len(portfolio_series) < 2):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "RISK ANALYSIS", font_size=22, bold=True, color=C_WHITE)

    y_pos = Inches(1.2)

    # Risk metrics table (left)
    if risk:
        filtered = {k: v for k, v in risk.items() if v is not None and isinstance(v, (int, float))}
        n_rows = len(filtered) + 1
        table_shape = slide.shapes.add_table(
            n_rows, 2,
            Inches(0.5), y_pos,
            Inches(4.2), Inches(0.35) * n_rows,
        )
        table = table_shape.table
        _set_cell_text(table.cell(0, 0), "Risk Metric", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
        _set_cell_text(table.cell(0, 1), "Value", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = C_NAVY
        table.cell(0, 1).fill.solid()
        table.cell(0, 1).fill.fore_color.rgb = C_NAVY

        for i, (metric, val) in enumerate(filtered.items()):
            _set_cell_text(table.cell(i + 1, 0), metric, font_size=9, alignment=PP_ALIGN.LEFT)
            if "sharpe" in metric.lower() or "ratio" in metric.lower() or "beta" in metric.lower() or "r-squared" in metric.lower():
                text = f"{val:.2f}"
            else:
                text = f"{val:.2%}"
            _set_cell_text(table.cell(i + 1, 1), text, font_size=9, alignment=PP_ALIGN.CENTER)
            if i % 2 == 0:
                table.cell(i + 1, 0).fill.solid()
                table.cell(i + 1, 0).fill.fore_color.rgb = C_LIGHT_BG
                table.cell(i + 1, 1).fill.solid()
                table.cell(i + 1, 1).fill.fore_color.rgb = C_LIGHT_BG

    # Daily returns statistics table (right)
    if portfolio_series is not None and len(portfolio_series) > 1:
        daily_rets = portfolio_series.pct_change().dropna()
        if len(daily_rets) > 0:
            stats = [
                ("Mean Daily Return", f"{daily_rets.mean():.4%}"),
                ("Std Dev (Daily)", f"{daily_rets.std():.4%}"),
                ("Skewness", f"{daily_rets.skew():.2f}"),
                ("Kurtosis", f"{daily_rets.kurtosis():.2f}"),
                ("Min Daily Return", f"{daily_rets.min():.4%}"),
                ("Max Daily Return", f"{daily_rets.max():.4%}"),
                ("Trading Days", f"{len(daily_rets):,}"),
            ]
            n_rows_stats = len(stats) + 1
            table_shape2 = slide.shapes.add_table(
                n_rows_stats, 2,
                Inches(5.2), y_pos,
                Inches(4.2), Inches(0.35) * n_rows_stats,
            )
            table2 = table_shape2.table
            _set_cell_text(table2.cell(0, 0), "Statistic", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
            _set_cell_text(table2.cell(0, 1), "Value", font_size=10, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
            table2.cell(0, 0).fill.solid()
            table2.cell(0, 0).fill.fore_color.rgb = C_NAVY
            table2.cell(0, 1).fill.solid()
            table2.cell(0, 1).fill.fore_color.rgb = C_NAVY

            for i, (label, val_str) in enumerate(stats):
                _set_cell_text(table2.cell(i + 1, 0), label, font_size=9, alignment=PP_ALIGN.LEFT)
                _set_cell_text(table2.cell(i + 1, 1), val_str, font_size=9, alignment=PP_ALIGN.CENTER)
                if i % 2 == 0:
                    table2.cell(i + 1, 0).fill.solid()
                    table2.cell(i + 1, 0).fill.fore_color.rgb = C_LIGHT_BG
                    table2.cell(i + 1, 1).fill.solid()
                    table2.cell(i + 1, 1).fill.fore_color.rgb = C_LIGHT_BG


def _slide_calendar_returns(prs: Presentation, report_data: dict):
    """Slide: Calendar year returns table."""
    cal_year_returns = report_data.get("calendar_year_returns", {})
    if not cal_year_returns:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "CALENDAR YEAR RETURNS", font_size=22, bold=True, color=C_WHITE)

    years = list(cal_year_returns.keys())
    returns = list(cal_year_returns.values())
    n_rows = len(years) + 1

    table_shape = slide.shapes.add_table(
        n_rows, 2,
        Inches(2.5), Inches(1.3),
        Inches(5.0), Inches(0.38) * n_rows,
    )
    table = table_shape.table

    _set_cell_text(table.cell(0, 0), "Year", font_size=11, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    _set_cell_text(table.cell(0, 1), "Return", font_size=11, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = C_NAVY
    table.cell(0, 1).fill.solid()
    table.cell(0, 1).fill.fore_color.rgb = C_NAVY

    for i, (yr, ret) in enumerate(zip(years, returns)):
        _set_cell_text(table.cell(i + 1, 0), str(yr), font_size=10, alignment=PP_ALIGN.CENTER)
        if isinstance(ret, (int, float)):
            color = C_GREEN if ret >= 0 else C_ACCENT
            _set_cell_text(table.cell(i + 1, 1), f"{ret:.2%}", font_size=10,
                           color=color, alignment=PP_ALIGN.CENTER)
        else:
            _set_cell_text(table.cell(i + 1, 1), str(ret), font_size=10, alignment=PP_ALIGN.CENTER)
        if i % 2 == 0:
            table.cell(i + 1, 0).fill.solid()
            table.cell(i + 1, 0).fill.fore_color.rgb = C_LIGHT_BG
            table.cell(i + 1, 1).fill.solid()
            table.cell(i + 1, 1).fill.fore_color.rgb = C_LIGHT_BG
    """Slide: Manager / AI commentary."""
    commentary = report_data.get("commentary")
    if not commentary:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = prs.slide_width

    # Header
    _add_shape_bg(slide, Emu(0), Emu(0), sw, Inches(0.9), C_DARK)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
                 "ANALYSIS & COMMENTARY", font_size=22, bold=True, color=C_WHITE)

    if isinstance(commentary, str):
        commentary = [commentary]

    y_pos = Inches(1.3)
    for para_text in commentary:
        tb = _add_textbox(slide, Inches(0.8), y_pos, Inches(8.4), Inches(1.0),
                          para_text, font_size=11, color=C_BODY)
        tb.text_frame.word_wrap = True
        y_pos += Inches(0.8)

    # Footer
    _add_textbox(slide, Inches(0.8), Inches(6.6), Inches(8), Inches(0.3),
                 "This report is for informational purposes only and does not constitute investment advice.",
                 font_size=8, color=C_MUTED, alignment=PP_ALIGN.LEFT)


# ──────────────────────────────────────────────
# Main entry point
# ──────────────────────────────────────────────
def generate_pptx_report(report_data: dict, output_path: str) -> str:
    """
    Generate a professional PowerPoint portfolio report.

    Parameters
    ----------
    report_data : dict
        Same schema as the PDF generator:
        - client_name, report_date, portfolio_value, currency
        - performance (dict), risk (dict), weights (dict)
        - holdings_df (DataFrame), portfolio_series, benchmark_series
        - portfolio_label, benchmark_label, commentary
    output_path : str
        Destination .pptx file path.

    Returns
    -------
    str : absolute path to the generated file.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, report_data)
    _slide_kpis(prs, report_data)
    _slide_allocation(prs, report_data)
    _slide_holdings(prs, report_data)
    _slide_performance(prs, report_data)
    _slide_frontier(prs, report_data)
    _slide_cumulative_returns(prs, report_data)
    _slide_correlation(prs, report_data)
    _slide_betas_capm(prs, report_data)
    _slide_risk_analysis(prs, report_data)
    _slide_calendar_returns(prs, report_data)
    _slide_commentary(prs, report_data)

    prs.save(output_path)
    return os.path.abspath(output_path)


# ──────────────────────────────────────────────
# Standalone test
# ──────────────────────────────────────────────
if __name__ == "__main__":
    np.random.seed(42)
    dates = pd.date_range("2024-01-01", periods=252, freq="B")
    cum = (1 + pd.Series(np.random.normal(0.0004, 0.012, len(dates)), index=dates)).cumprod() * 100
    bench = (1 + pd.Series(np.random.normal(0.0003, 0.011, len(dates)), index=dates)).cumprod() * 100

    holdings = pd.DataFrame({
        "Ticker": ["SPY", "AGG", "GLD", "VNQ"],
        "Weight": [0.35, 0.40, 0.15, 0.10],
        "Market Value": [87500.0, 100000.0, 37500.0, 25000.0],
        "P/L": [4200.0, 1500.0, 2100.0, -300.0],
    })

    data = {
        "client_name": "Jane Doe",
        "report_date": "2026-02-06",
        "portfolio_value": 250000,
        "currency": "EUR",
        "performance": {"1M": 0.012, "YTD": 0.034, "1Y": 0.091, "Since Inception": 0.184},
        "risk": {"Volatility": 0.145, "Sharpe": 1.12, "Max Drawdown": -0.087},
        "weights": {"SPY": 0.35, "AGG": 0.40, "GLD": 0.15, "VNQ": 0.10},
        "holdings_df": holdings,
        "portfolio_series": cum,
        "benchmark_series": bench,
        "portfolio_label": "Max Sharpe",
        "benchmark_label": "S&P 500",
        "commentary": [
            "The portfolio maintained a defensive tilt with 40% allocation to bonds.",
            "Equity exposure via broad market ETFs captured the late-year rally.",
        ],
    }

    out = generate_pptx_report(data, "test_portfolio_report.pptx")
    print(f"Report saved to: {out}")
